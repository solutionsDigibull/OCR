"""PowerPoint file processor."""

import os
import logging
from typing import Dict, Any

from .base import BaseProcessor
from ..result import ConversionResult
from ..exceptions import ConversionError, FileNotFoundError

# Configure logging
logger = logging.getLogger(__name__)


class PPTXProcessor(BaseProcessor):
    """Processor for PowerPoint files (PPT, PPTX)."""
    
    def can_process(self, file_path: str) -> bool:
        """Check if this processor can handle the given file.
        
        Args:
            file_path: Path to the file to check
            
        Returns:
            True if this processor can handle the file
        """
        if not os.path.exists(file_path):
            return False
        
        # Check file extension - ensure file_path is a string
        file_path_str = str(file_path)
        _, ext = os.path.splitext(file_path_str.lower())
        return ext in ['.ppt', '.pptx']
    
    def process(self, file_path: str) -> ConversionResult:
        """Process the PowerPoint file and return a conversion result.
        
        Args:
            file_path: Path to the PowerPoint file to process
            
        Returns:
            ConversionResult containing the processed content
            
        Raises:
            FileNotFoundError: If the file doesn't exist
            ConversionError: If processing fails
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Initialize metadata
        metadata = {
            "file_path": file_path,
            "file_size": os.path.getsize(file_path),
            "processor": "PPTXProcessor"
        }
        
        # Check file extension to determine processing method
        file_path_str = str(file_path)
        _, ext = os.path.splitext(file_path_str.lower())
        
        if ext == '.ppt':
            return self._process_ppt_file(file_path, metadata)
        else:
            return self._process_pptx_file(file_path, metadata)
    
    def _process_ppt_file(self, file_path: str, metadata: Dict[str, Any]) -> ConversionResult:
        """Process .ppt files using pypandoc."""
        try:
            import pypandoc
            
            # Convert .ppt to markdown using pandoc
            content = pypandoc.convert_file(file_path, 'markdown')
            
            metadata.update({
                "file_type": "ppt",
                "extractor": "pypandoc"
            })
            
            # Clean up the content
            content = self._clean_content(content)
            
            return ConversionResult(content, metadata)
            
        except ImportError:
            raise ConversionError("pypandoc is required for .ppt file processing. Install it with: pip install pypandoc")
        except Exception as e:
            raise ConversionError(f"Failed to process .ppt file {file_path}: {str(e)}")
    
    def _process_pptx_file(self, file_path: str, metadata: Dict[str, Any]) -> ConversionResult:
        """Process .pptx files using python-pptx."""
        try:
            from pptx import Presentation
            
            content_parts = []
            prs = Presentation(file_path)
            
            metadata.update({
                "slide_count": len(prs.slides),
                "file_type": "pptx",
                "extractor": "python-pptx"
            })
            
            # Check if preserve_layout is available (from base class or config)
            preserve_layout = getattr(self, 'preserve_layout', False)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                if preserve_layout:
                    content_parts.append(f"\n## Slide {slide_num}\n")
                
                slide_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    content_parts.extend(slide_content)
                    content_parts.append("")  # Add spacing between slides
            
            content = "\n\n".join(content_parts)
            
            # Clean up the content
            content = self._clean_content(content)
            
            return ConversionResult(content, metadata)
            
        except ImportError:
            raise ConversionError("python-pptx is required for .pptx file processing. Install it with: pip install python-pptx")
        except Exception as e:
            if isinstance(e, (FileNotFoundError, ConversionError)):
                raise
            raise ConversionError(f"Failed to process .pptx file {file_path}: {str(e)}")
    
    def _clean_content(self, content: str) -> str:
        """Clean up the extracted PowerPoint content.
        
        Args:
            content: Raw PowerPoint text content
            
        Returns:
            Cleaned text content
        """
        # Remove excessive whitespace and normalize
        lines = content.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # Remove excessive whitespace
            line = ' '.join(line.split())
            if line.strip():
                cleaned_lines.append(line)
        
        # Join lines and add proper spacing
        content = '\n'.join(cleaned_lines)
        
        # Add spacing around headers
        content = content.replace('## Slide', '\n## Slide')
        
        return content.strip() 